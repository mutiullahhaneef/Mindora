"""
Document Parser — extracts raw text from PDF, PPTX, DOCX, and TXT files.
"""
import logging

logger = logging.getLogger(__name__)


class DocumentParser:
    """Routes file parsing to the correct format-specific extractor."""

    def parse(self, file_path: str, file_type: str) -> str:
        """
        Extract text from a file.
        Args:
            file_path: absolute or relative path to the file on disk
            file_type: one of "pdf", "pptx", "docx", "txt"
        Returns:
            Extracted text as a single string.
        Raises:
            ValueError: if file_type is not supported.
        """
        parsers = {
            "pdf": self._parse_pdf,
            "pptx": self._parse_pptx,
            "docx": self._parse_docx,
            "txt": self._parse_txt,
        }
        parser_fn = parsers.get(file_type.lower())
        if not parser_fn:
            raise ValueError(f"Unsupported file type for parsing: {file_type}")

        try:
            text = parser_fn(file_path)
            logger.info("Parsed %s (%s): %d chars", file_path, file_type, len(text))
            return text
        except Exception as exc:
            logger.error("Failed to parse %s: %s", file_path, exc)
            raise

    # ─── Format-specific parsers ───────────────────────────────────────────

    def _parse_pdf(self, path: str) -> str:
        """Use PyMuPDF (fitz) to extract text from all pages."""
        import fitz  # PyMuPDF

        doc = fitz.open(path)
        pages: list[str] = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n".join(pages)

    def _parse_pptx(self, path: str) -> str:
        """Use python-pptx to iterate slides and extract all text frames."""
        from pptx import Presentation

        prs = Presentation(path)
        slides: list[str] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_texts.append(line)
            slides.append("\n".join(slide_texts))
        return "\n\n".join(slides)

    def _parse_docx(self, path: str) -> str:
        """Use python-docx to extract all paragraph text."""
        from docx import Document

        doc = Document(path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(paragraphs)

    def _parse_txt(self, path: str) -> str:
        """Read plain text with UTF-8, falling back to latin-1."""
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1") as f:
                return f.read()
